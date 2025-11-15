"""
Document Processing Integration
Handles PDF, Word, Excel, PowerPoint, and other document formats
"""

from typing import Dict, Any, List, Optional, BinaryIO
import logging
import os
import io
from datetime import datetime
from pathlib import Path

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """
    Document processor for multiple file formats
    """
    
    def __init__(self):
        self.supported_formats = [
            'pdf', 'docx', 'doc', 'xlsx', 'xls', 
            'pptx', 'ppt', 'txt', 'md', 'csv', 'html'
        ]
        self.image_formats = ['jpg', 'jpeg', 'png', 'tiff', 'bmp']
    
    async def extract_text(self, file_path: str, file_type: str) -> str:
        """
        Extract text from document
        
        Args:
            file_path: Path to document
            file_type: Document type (pdf, docx, etc.)
            
        Returns:
            Extracted text
        """
        try:
            file_type = file_type.lower().replace('.', '')
            
            if file_type == 'pdf':
                return await self._extract_text_from_pdf(file_path)
            elif file_type in ['docx', 'doc']:
                return await self._extract_text_from_word(file_path)
            elif file_type in ['xlsx', 'xls']:
                return await self._extract_text_from_excel(file_path)
            elif file_type in ['pptx', 'ppt']:
                return await self._extract_text_from_powerpoint(file_path)
            elif file_type in ['txt', 'md', 'csv']:
                return await self._extract_text_from_text(file_path)
            elif file_type == 'html':
                return await self._extract_text_from_html(file_path)
            elif file_type in self.image_formats:
                return await self.perform_ocr(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
                
        except Exception as e:
            logger.error(f"Error extracting text: {str(e)}")
            raise
    
    async def _extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF using pdfplumber"""
        try:
            import pdfplumber
            
            text_parts = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
            
            return '\n\n'.join(text_parts)
            
        except ImportError:
            # Fallback to PyPDF2
            try:
                import PyPDF2
                
                text_parts = []
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        text = page.extract_text()
                        if text:
                            text_parts.append(text)
                
                return '\n\n'.join(text_parts)
            except Exception as e:
                logger.error(f"Error with PyPDF2: {str(e)}")
                raise
        except Exception as e:
            logger.error(f"Error extracting PDF text: {str(e)}")
            raise
    
    async def _extract_text_from_word(self, file_path: str) -> str:
        """Extract text from Word document using python-docx"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            text_parts = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = '\t'.join(cell.text for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error extracting Word text: {str(e)}")
            raise
    
    async def _extract_text_from_excel(self, file_path: str) -> str:
        """Extract text from Excel using openpyxl"""
        try:
            from openpyxl import load_workbook
            
            workbook = load_workbook(file_path, data_only=True)
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"Sheet: {sheet_name}")
                text_parts.append("-" * 50)
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                    if row_text.strip():
                        text_parts.append(row_text)
                
                text_parts.append("")
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error extracting Excel text: {str(e)}")
            raise
    
    async def _extract_text_from_powerpoint(self, file_path: str) -> str:
        """Extract text from PowerPoint using python-pptx"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_parts = []
            
            for i, slide in enumerate(prs.slides, 1):
                text_parts.append(f"Slide {i}:")
                text_parts.append("-" * 50)
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
                
                text_parts.append("")
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error extracting PowerPoint text: {str(e)}")
            raise
    
    async def _extract_text_from_text(self, file_path: str) -> str:
        """Extract text from text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except:
                    continue
            raise ValueError("Could not decode file with any known encoding")
    
    async def _extract_text_from_html(self, file_path: str) -> str:
        """Extract text from HTML"""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                
                # Get text
                text = soup.get_text()
                
                # Clean up whitespace
                lines = (line.strip() for line in text.splitlines())
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                text = '\n'.join(chunk for chunk in chunks if chunk)
                
                return text
                
        except Exception as e:
            logger.error(f"Error extracting HTML text: {str(e)}")
            # Fallback to plain text extraction
            return await self._extract_text_from_text(file_path)
    
    async def extract_tables(self, file_path: str, file_type: str) -> List[Dict]:
        """
        Extract tables from document
        
        Args:
            file_path: Path to document
            file_type: Document type
            
        Returns:
            List of tables as dictionaries
        """
        try:
            file_type = file_type.lower().replace('.', '')
            
            if file_type == 'pdf':
                return await self._extract_tables_from_pdf(file_path)
            elif file_type in ['docx', 'doc']:
                return await self._extract_tables_from_word(file_path)
            elif file_type in ['xlsx', 'xls']:
                return await self._extract_tables_from_excel(file_path)
            elif file_type == 'html':
                return await self._extract_tables_from_html(file_path)
            else:
                return []
                
        except Exception as e:
            logger.error(f"Error extracting tables: {str(e)}")
            raise
    
    async def _extract_tables_from_pdf(self, file_path: str) -> List[Dict]:
        """Extract tables from PDF"""
        try:
            import pdfplumber
            
            tables = []
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    page_tables = page.extract_tables()
                    for table_num, table in enumerate(page_tables, 1):
                        if table:
                            tables.append({
                                'page': page_num,
                                'table_number': table_num,
                                'data': table,
                                'rows': len(table),
                                'columns': len(table[0]) if table else 0
                            })
            
            return tables
            
        except Exception as e:
            logger.error(f"Error extracting PDF tables: {str(e)}")
            return []
    
    async def _extract_tables_from_word(self, file_path: str) -> List[Dict]:
        """Extract tables from Word document"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            tables = []
            
            for table_num, table in enumerate(doc.tables, 1):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                
                tables.append({
                    'table_number': table_num,
                    'data': table_data,
                    'rows': len(table_data),
                    'columns': len(table_data[0]) if table_data else 0
                })
            
            return tables
            
        except Exception as e:
            logger.error(f"Error extracting Word tables: {str(e)}")
            return []
    
    async def _extract_tables_from_excel(self, file_path: str) -> List[Dict]:
        """Extract tables from Excel"""
        try:
            from openpyxl import load_workbook
            
            workbook = load_workbook(file_path, data_only=True)
            tables = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # Get all data
                data = []
                for row in sheet.iter_rows(values_only=True):
                    data.append(list(row))
                
                if data:
                    tables.append({
                        'sheet': sheet_name,
                        'data': data,
                        'rows': len(data),
                        'columns': len(data[0]) if data else 0
                    })
            
            return tables
            
        except Exception as e:
            logger.error(f"Error extracting Excel tables: {str(e)}")
            return []
    
    async def _extract_tables_from_html(self, file_path: str) -> List[Dict]:
        """Extract tables from HTML"""
        try:
            import pandas as pd
            
            tables_df = pd.read_html(file_path)
            tables = []
            
            for i, df in enumerate(tables_df, 1):
                tables.append({
                    'table_number': i,
                    'data': df.values.tolist(),
                    'headers': df.columns.tolist(),
                    'rows': len(df),
                    'columns': len(df.columns)
                })
            
            return tables
            
        except Exception as e:
            logger.error(f"Error extracting HTML tables: {str(e)}")
            return []
    
    async def extract_images(self, file_path: str, file_type: str, output_dir: str) -> List[str]:
        """
        Extract images from document
        
        Args:
            file_path: Path to document
            file_type: Document type
            output_dir: Directory to save extracted images
            
        Returns:
            List of image file paths
        """
        try:
            file_type = file_type.lower().replace('.', '')
            
            if file_type == 'pdf':
                return await self._extract_images_from_pdf(file_path, output_dir)
            elif file_type in ['docx', 'doc']:
                return await self._extract_images_from_word(file_path, output_dir)
            elif file_type in ['pptx', 'ppt']:
                return await self._extract_images_from_powerpoint(file_path, output_dir)
            else:
                return []
                
        except Exception as e:
            logger.error(f"Error extracting images: {str(e)}")
            raise
    
    async def _extract_images_from_pdf(self, file_path: str, output_dir: str) -> List[str]:
        """Extract images from PDF"""
        try:
            import fitz  # PyMuPDF
            
            os.makedirs(output_dir, exist_ok=True)
            image_paths = []
            
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc[page_num]
                images = page.get_images()
                
                for img_num, img in enumerate(images):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    image_path = os.path.join(output_dir, f"page{page_num+1}_img{img_num+1}.{image_ext}")
                    with open(image_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    
                    image_paths.append(image_path)
            
            return image_paths
            
        except Exception as e:
            logger.error(f"Error extracting PDF images: {str(e)}")
            return []
    
    async def _extract_images_from_word(self, file_path: str, output_dir: str) -> List[str]:
        """Extract images from Word document"""
        try:
            from docx import Document
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
            
            os.makedirs(output_dir, exist_ok=True)
            image_paths = []
            
            doc = Document(file_path)
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image = rel.target_part.blob
                    image_ext = rel.target_ref.split('.')[-1]
                    image_path = os.path.join(output_dir, f"image_{len(image_paths)+1}.{image_ext}")
                    
                    with open(image_path, 'wb') as img_file:
                        img_file.write(image)
                    
                    image_paths.append(image_path)
            
            return image_paths
            
        except Exception as e:
            logger.error(f"Error extracting Word images: {str(e)}")
            return []
    
    async def _extract_images_from_powerpoint(self, file_path: str, output_dir: str) -> List[str]:
        """Extract images from PowerPoint"""
        try:
            from pptx import Presentation
            
            os.makedirs(output_dir, exist_ok=True)
            image_paths = []
            
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, 1):
                for shape_num, shape in enumerate(slide.shapes, 1):
                    if shape.shape_type == 13:  # Picture
                        image = shape.image
                        image_bytes = image.blob
                        image_ext = image.ext
                        
                        image_path = os.path.join(output_dir, f"slide{slide_num}_img{shape_num}.{image_ext}")
                        with open(image_path, 'wb') as img_file:
                            img_file.write(image_bytes)
                        
                        image_paths.append(image_path)
            
            return image_paths
            
        except Exception as e:
            logger.error(f"Error extracting PowerPoint images: {str(e)}")
            return []
    
    async def extract_metadata(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """
        Extract metadata from document
        
        Args:
            file_path: Path to document
            file_type: Document type
            
        Returns:
            Metadata dictionary
        """
        try:
            file_type = file_type.lower().replace('.', '')
            
            # Get file stats
            file_stats = os.stat(file_path)
            base_metadata = {
                "filename": os.path.basename(file_path),
                "file_size": file_stats.st_size,
                "created_date": datetime.fromtimestamp(file_stats.st_ctime).isoformat(),
                "modified_date": datetime.fromtimestamp(file_stats.st_mtime).isoformat(),
            }
            
            if file_type == 'pdf':
                metadata = await self._extract_metadata_from_pdf(file_path)
            elif file_type in ['docx', 'doc']:
                metadata = await self._extract_metadata_from_word(file_path)
            elif file_type in ['xlsx', 'xls']:
                metadata = await self._extract_metadata_from_excel(file_path)
            elif file_type in ['pptx', 'ppt']:
                metadata = await self._extract_metadata_from_powerpoint(file_path)
            else:
                metadata = {}
            
            return {**base_metadata, **metadata}
            
        except Exception as e:
            logger.error(f"Error extracting metadata: {str(e)}")
            raise
    
    async def _extract_metadata_from_pdf(self, file_path: str) -> Dict[str, Any]:
        """Extract metadata from PDF"""
        try:
            import PyPDF2
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                metadata = pdf_reader.metadata
                
                return {
                    "author": metadata.get('/Author', 'Unknown'),
                    "title": metadata.get('/Title', ''),
                    "subject": metadata.get('/Subject', ''),
                    "creator": metadata.get('/Creator', ''),
                    "producer": metadata.get('/Producer', ''),
                    "page_count": len(pdf_reader.pages)
                }
                
        except Exception as e:
            logger.error(f"Error extracting PDF metadata: {str(e)}")
            return {}
    
    async def _extract_metadata_from_word(self, file_path: str) -> Dict[str, Any]:
        """Extract metadata from Word document"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            core_props = doc.core_properties
            
            return {
                "author": core_props.author or 'Unknown',
                "title": core_props.title or '',
                "subject": core_props.subject or '',
                "keywords": core_props.keywords or '',
                "created": core_props.created.isoformat() if core_props.created else None,
                "modified": core_props.modified.isoformat() if core_props.modified else None,
                "page_count": len(doc.sections),
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables)
            }
            
        except Exception as e:
            logger.error(f"Error extracting Word metadata: {str(e)}")
            return {}
    
    async def _extract_metadata_from_excel(self, file_path: str) -> Dict[str, Any]:
        """Extract metadata from Excel"""
        try:
            from openpyxl import load_workbook
            
            workbook = load_workbook(file_path)
            props = workbook.properties
            
            return {
                "author": props.creator or 'Unknown',
                "title": props.title or '',
                "subject": props.subject or '',
                "keywords": props.keywords or '',
                "created": props.created.isoformat() if props.created else None,
                "modified": props.modified.isoformat() if props.modified else None,
                "sheet_count": len(workbook.sheetnames),
                "sheet_names": workbook.sheetnames
            }
            
        except Exception as e:
            logger.error(f"Error extracting Excel metadata: {str(e)}")
            return {}
    
    async def _extract_metadata_from_powerpoint(self, file_path: str) -> Dict[str, Any]:
        """Extract metadata from PowerPoint"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            core_props = prs.core_properties
            
            return {
                "author": core_props.author or 'Unknown',
                "title": core_props.title or '',
                "subject": core_props.subject or '',
                "keywords": core_props.keywords or '',
                "created": core_props.created.isoformat() if core_props.created else None,
                "modified": core_props.modified.isoformat() if core_props.modified else None,
                "slide_count": len(prs.slides)
            }
            
        except Exception as e:
            logger.error(f"Error extracting PowerPoint metadata: {str(e)}")
            return {}
    
    async def perform_ocr(self, file_path: str, language: str = "eng") -> str:
        """
        Perform OCR on document or image
        
        Args:
            file_path: Path to document or image
            language: OCR language (default: eng)
            
        Returns:
            Extracted text
        """
        try:
            import pytesseract
            from PIL import Image
            
            # Open image
            image = Image.open(file_path)
            
            # Perform OCR
            text = pytesseract.image_to_string(image, lang=language)
            
            return text
            
        except Exception as e:
            logger.error(f"Error performing OCR: {str(e)}")
            raise
    
    async def convert_document(
        self, 
        file_path: str, 
        source_format: str, 
        target_format: str,
        output_path: Optional[str] = None
    ) -> str:
        """
        Convert document to different format
        
        Args:
            file_path: Path to source document
            source_format: Source format
            target_format: Target format
            output_path: Optional output path
            
        Returns:
            Path to converted document
        """
        try:
            source_format = source_format.lower().replace('.', '')
            target_format = target_format.lower().replace('.', '')
            
            if not output_path:
                base_name = os.path.splitext(file_path)[0]
                output_path = f"{base_name}.{target_format}"
            
            # PDF conversions
            if source_format == 'pdf' and target_format in ['txt', 'text']:
                text = await self._extract_text_from_pdf(file_path)
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(text)
                return output_path
            
            # Add more conversion logic as needed
            raise NotImplementedError(f"Conversion from {source_format} to {target_format} not yet implemented")
            
        except Exception as e:
            logger.error(f"Error converting document: {str(e)}")
            raise
    
    async def search_document(self, file_path: str, file_type: str, query: str) -> List[Dict]:
        """
        Search within document
        
        Args:
            file_path: Path to document
            file_type: Document type
            query: Search query
            
        Returns:
            List of search results with page numbers and context
        """
        try:
            # Extract text
            text = await self.extract_text(file_path, file_type)
            
            # Search for query
            results = []
            lines = text.split('\n')
            query_lower = query.lower()
            
            for line_num, line in enumerate(lines, 1):
                if query_lower in line.lower():
                    # Get context (previous and next lines)
                    start = max(0, line_num - 2)
                    end = min(len(lines), line_num + 1)
                    context = '\n'.join(lines[start:end])
                    
                    results.append({
                        'line_number': line_num,
                        'match': line.strip(),
                        'context': context,
                        'position': line.lower().index(query_lower)
                    })
            
            return results
            
        except Exception as e:
            logger.error(f"Error searching document: {str(e)}")
            raise
    
    async def compare_documents(self, file_path1: str, file_path2: str, file_type1: str, file_type2: str) -> Dict[str, Any]:
        """
        Compare two documents
        
        Args:
            file_path1: Path to first document
            file_path2: Path to second document
            file_type1: First document type
            file_type2: Second document type
            
        Returns:
            Comparison results with differences and similarity score
        """
        try:
            # Extract text from both documents
            text1 = await self.extract_text(file_path1, file_type1)
            text2 = await self.extract_text(file_path2, file_type2)
            
            # Split into lines
            lines1 = text1.split('\n')
            lines2 = text2.split('\n')
            
            # Simple diff
            import difflib
            diff = list(difflib.unified_diff(lines1, lines2, lineterm=''))
            
            # Calculate similarity
            similarity = difflib.SequenceMatcher(None, text1, text2).ratio()
            
            # Categorize changes
            added = [line[1:] for line in diff if line.startswith('+') and not line.startswith('+++')]
            removed = [line[1:] for line in diff if line.startswith('-') and not line.startswith('---')]
            
            return {
                "similarity": round(similarity, 4),
                "differences": diff,
                "added": added,
                "removed": removed,
                "added_count": len(added),
                "removed_count": len(removed),
                "total_changes": len(added) + len(removed)
            }
            
        except Exception as e:
            logger.error(f"Error comparing documents: {str(e)}")
            raise